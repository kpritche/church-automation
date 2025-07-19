# ppt_generator_services.py
"""
Generate a PowerPoint (.pptx) per item, with:
- Black background blank slides
- Stateful coloring: white for calls, golden for responses, persisting until reset
- Prepend scripture reference slide if provided
- Song lyrics slices assumed to be prepared upstream (two lines per slide)
- 16:9 ratio at 1920×1080 pixels

Usage:
    from ppt_generator_services import create_pptx_for_items
    create_pptx_for_items(
        slides,                  # List[Dict[str,str]] slides to render
        filename,                # output .pptx path
        scripture_reference=None # Optional[str] reference slide text
    )
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from typing import List, Dict, Optional

# Speaker markers
CALL_MARKERS = ("Leader:", "L:")
RESPONSE_MARKERS = ("People:", "P:")
# Colors
WHITE = RGBColor(255, 255, 255)
RESPONSE_COLOR = RGBColor(255, 215, 0)  # lighter goldenrod


def create_pptx_for_items(
    slides: List[Dict[str, str]],
    filename: str = "service_slides.pptx",
    scripture_reference: Optional[str] = None
) -> None:
    """
    Render slides into a .pptx file.

    slides: each dict with 'text' and 'style' keys ('content' or 'blank').
    scripture_reference: text for an initial slide if provided.
    """
    # Initialize presentation
    prs = Presentation()
    # Set slide dimensions to 1920x1080 px at 96 DPI
    emu_per_px = int(914400 / 96)
    prs.slide_width = Emu(1920 * emu_per_px)
    prs.slide_height = Emu(1080 * emu_per_px)
    blank_layout = prs.slide_layouts[6]

    def add_text_slide(text: str, color: RGBColor) -> None:
        # Add a content slide with given text and text color
        sl = prs.slides.add_slide(blank_layout)
        # Black background
        bg = sl.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0, 0, 0)
        # Full-slide textbox
        tx = sl.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        tf = tx.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(80)
        run.font.name = 'Arial'
        run.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # Start in call (white) mode
    current_color = WHITE

    # Prepend scripture reference if present
    if scripture_reference:
        add_text_slide(scripture_reference, WHITE)

    for slide in slides:
        text = slide.get('text', '').strip()
        # Blank slide: always black background
        if slide.get('style') == 'blank' or not text:
            sl = prs.slides.add_slide(blank_layout)
            sl.background.fill.solid()
            sl.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
            continue

        # Update color state based on markers
        if any(text.startswith(m) for m in CALL_MARKERS):
            current_color = WHITE
        elif any(text.startswith(m) for m in RESPONSE_MARKERS):
            current_color = RESPONSE_COLOR

        # Render content slide
        add_text_slide(text, current_color)

    # Save the .pptx
    prs.save(filename)
