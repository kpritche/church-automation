import os
import sys
import uuid
import zipfile
import base64
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from google.protobuf.message import Message

# Add the generated protobuf modules to the path FIRST
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'ProPresenter7_Proto', 'generated'))

# Import protobufs
import presentation_pb2
import cue_pb2
import action_pb2
import slide_pb2
import graphicsData_pb2
import color_pb2
import uuid_pb2
import hotKey_pb2
import url_pb2
import background_pb2

class PptxToProConverter:
    def __init__(self, pptx_path, output_dir="output"):
        self.pptx_path = pptx_path
        self.output_dir = output_dir
        self.presentation_name = os.path.splitext(os.path.basename(pptx_path))[0]
        self.media_assets = {}  # Map filename -> bytes
        
        self.TARGET_WIDTH = 1920
        self.TARGET_HEIGHT = 1080
        
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def generate_uuid(self):
        u = uuid_pb2.UUID()
        u.string = str(uuid.uuid4())
        return u

    def create_color(self, red, green, blue, alpha=1.0):
        c = color_pb2.Color()
        c.red = red
        c.green = green
        c.blue = blue
        c.alpha = alpha
        return c

    def get_rgb_from_pptx_color(self, color_format, default=(255, 255, 255)):
        if not color_format:
            return default
        try:
            if hasattr(color_format, 'type'):
                if color_format.type == 1: # RGB
                    return (color_format.rgb[0], color_format.rgb[1], color_format.rgb[2])
                elif color_format.type == 2: # Theme color
                    # Fallback for theme colors
                    return default 
        except:
            pass
        return default

    def convert_position(self, emu_value, axis='x', pptx_size=None):
        if pptx_size is None or pptx_size == 0:
            return 0
        if axis == 'x':
            scale = self.TARGET_WIDTH / pptx_size
        else:
            scale = self.TARGET_HEIGHT / pptx_size
        return emu_value * scale

    def generate_rtf(self, text_frame):
        colors = []
        def get_color_index(rgb_tuple):
            if rgb_tuple not in colors:
                colors.append(rgb_tuple)
            return colors.index(rgb_tuple) + 1

        default_color = (255, 255, 255)
        rtf_content = ""
        
        for paragraph in text_frame.paragraphs:
            align_tag = r"\qc"
            if paragraph.alignment == PP_ALIGN.LEFT: align_tag = r"\ql"
            elif paragraph.alignment == PP_ALIGN.RIGHT: align_tag = r"\qr"
            elif paragraph.alignment == PP_ALIGN.JUSTIFY: align_tag = r"\qj"
            
            rtf_content += f"{align_tag} "
            
            for run in paragraph.runs:
                rgb = default_color
                if run.font.color:
                    rgb = self.get_rgb_from_pptx_color(run.font.color, default_color)
                
                color_idx = get_color_index(rgb)
                
                font_size = 100 
                if run.font.size:
                    font_size = int(run.font.size.pt * 2) 
                
                text = run.text.replace('\\', '\\\\').replace('{', '\{').replace('}', '\}').replace('\n', '\\\n')
                
                style_tags = ""
                if run.font.bold: style_tags += r"\b"
                if run.font.italic: style_tags += r"\i"
                if run.font.underline: style_tags += r"\ul"
                
                close_tags = ""
                if run.font.bold: close_tags += r"\b0"
                if run.font.italic: close_tags += r"\i0"
                if run.font.underline: close_tags += r"\ulnone"

                rtf_content += f"{style_tags}\\cf{color_idx}\\fs{font_size} {text}{close_tags}"
            
            rtf_content += r"\par "

        colortbl = r"{\colortbl;"
        for r, g, b in colors:
            colortbl += f"\\red{r}\\green{g}\\blue{b};"
        colortbl += "}"

        full_rtf = (
            r"{\rtf1\ansi\ansicpg1252\cocoartf1671\cocoasubrtf600"
            r"{\fonttbl\f0\fswiss\fcharset0 Helvetica;}"
            f"{colortbl}"
            r"\pard\pardirnatural\partightenfactor0"
            f"\f0 {rtf_content}"
            r"}"
        )
        return full_rtf.encode('utf-8')

    def convert(self):
        print(f"Reading {self.pptx_path}...")
        prs = PptxPresentation(self.pptx_path)
        
        pptx_width = prs.slide_width
        pptx_height = prs.slide_height
        
        pro_pres = presentation_pb2.Presentation()
        pro_pres.uuid.CopyFrom(self.generate_uuid())
        pro_pres.name = self.presentation_name
        pro_pres.selected_arrangement.CopyFrom(self.generate_uuid())
        
        for i, slide in enumerate(prs.slides):
            print(f"Processing slide {i+1}...")
            
            cue = pro_pres.cues.add()
            cue.uuid.CopyFrom(self.generate_uuid())
            cue.name = f"Slide {i+1}"
            cue.completion_target_uuid.CopyFrom(self.generate_uuid())
            cue.completion_action_uuid.CopyFrom(self.generate_uuid())
            
            action = cue.actions.add()
            action.uuid.CopyFrom(self.generate_uuid())
            
            # Safe Action Type Assignment
            try:
                action.type = action_pb2.Action.ACTION_TYPE_PRESENTATION_SLIDE
            except AttributeError:
                # Fallback to integer if enum is unresolved
                action.type = 0
            
            pro_slide = action.slide.presentation.base_slide
            pro_slide.uuid.CopyFrom(self.generate_uuid())
            pro_slide.size.width = self.TARGET_WIDTH
            pro_slide.size.height = self.TARGET_HEIGHT
            
            # Background Color Logic
            bg_color = (0, 0, 0)
            try:
                bg = slide.background
                if bg and bg.fill.type == 1:
                    c = bg.fill.fore_color
                    bg_color = self.get_rgb_from_pptx_color(c, (0,0,0))
            except:
                pass
            
            # Try setting background color - Attempt 1: Direct Field
            success_bg = False
            try:
                pro_slide.background_color.CopyFrom(
                    self.create_color(bg_color[0]/255.0, bg_color[1]/255.0, bg_color[2]/255.0, 1.0)
                )
                pro_slide.draws_background_color = True
                success_bg = True
            except (AttributeError, ValueError):
                pass
            
            # Try setting background color - Attempt 2: Nested Background Object
            if not success_bg:
                try:
                    pro_slide.background.color.CopyFrom(
                        self.create_color(bg_color[0]/255.0, bg_color[1]/255.0, bg_color[2]/255.0, 1.0)
                    )
                    pro_slide.background.enable = True
                except (AttributeError, ValueError):
                    # If both fail, we just skip background color to avoid crashing
                    # print(f"Warning: Could not set background color structure for slide {i+1}")
                    pass

            for shape in slide.shapes:
                try:
                    x = self.convert_position(shape.left, 'x', pptx_width)
                    y = self.convert_position(shape.top, 'y', pptx_height)
                    w = self.convert_position(shape.width, 'x', pptx_width)
                    h = self.convert_position(shape.height, 'y', pptx_height)
                except AttributeError:
                    continue 

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_bytes = shape.image.blob
                        ext = shape.image.ext
                        img_name = f"img_{uuid.uuid4().hex[:8]}.{ext}"
                        self.media_assets[img_name] = image_bytes
                        
                        element_container = pro_slide.elements.add()
                        element = element_container.element
                        element.uuid.CopyFrom(self.generate_uuid())
                        element.name = "Image"
                        element.bounds.origin.x = x
                        element.bounds.origin.y = y
                        element.bounds.size.width = w
                        element.bounds.size.height = h
                        element.opacity = 1.0
                        
                        url_obj = url_pb2.URL()
                        url_obj.string = img_name
                        
                        element.fill.enable = True
                        element.fill.media.url.CopyFrom(url_obj)
                        
                    except Exception as e:
                        print(f"Failed to process image on slide {i+1}: {e}")

                elif shape.has_text_frame and shape.text.strip():
                    element_container = pro_slide.elements.add()
                    element = element_container.element
                    element.uuid.CopyFrom(self.generate_uuid())
                    element.name = "Text Box"
                    
                    element.bounds.origin.x = x
                    element.bounds.origin.y = y
                    element.bounds.size.width = w
                    element.bounds.size.height = h
                    element.opacity = 1.0
                    
                    element.text.rtf_data = self.generate_rtf(shape.text_frame)

        pro_filename = f"{self.presentation_name}.pro"
        pro_path = os.path.join(self.output_dir, pro_filename)
        
        with open(pro_path, "wb") as f:
            f.write(pro_pres.SerializeToString())
            
        print(f"Created .pro file at {pro_path}")
        
        bundle_name = f"{self.presentation_name}.proBundle"
        bundle_path = os.path.join(self.output_dir, bundle_name)
        
        with zipfile.ZipFile(bundle_path, 'w') as zipf:
            zipf.write(pro_path, arcname=pro_filename)
            for filename, data in self.media_assets.items():
                zipf.writestr(filename, data)
                print(f"Included asset: {filename}")
            
        print(f"Successfully created bundle: {bundle_path}")
        return bundle_path