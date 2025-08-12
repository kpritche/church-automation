import qrcode
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import requests
import config
from PIL import Image
import os
from win32com.client import Dispatch
import textwrap
from summarize import summarize_title

def truncate_text(text, max_length):
    """Truncate text to a maximum length, ensuring it ends with a complete word."""
    if len(text) <= max_length:
        return text
    truncated = text[:max_length].rsplit(' ', 1)[0]
    return truncated if truncated else text[:max_length] + '...'

def determine_font_size(text, base_size, min_size):
    """Determine the font size based on text length."""
    if len(text) <= 500:
        return base_size
    elif len(text) <= 600:
        return base_size - 2
    elif len(text) <= 700:
        return base_size - 4
    elif len(text) <= 800:
        return base_size - 6
    else:
        return max(min_size, base_size - 8)

def generate_qr_code(link):
    qr = qrcode.QRCode(border=1)
    qr.add_data(link)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buffer = BytesIO()
    img.save(buffer, format='PNG')
    buffer.seek(0)
    return buffer

def fetch_image_from_url(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception as e:
        print(f"Error fetching image from {url}: {e}")
        return None
    
def export_pptx_to_jpg(pptx_path: str, output_dir: str):
    """
    Opens the given PPTX in PowerPoint via COM and exports each slide
    as a JPG into output_dir (creates it if necessary).
    """
    # Ensure output folder exists
    os.makedirs(output_dir, exist_ok=True)

    # Launch PowerPoint
    ppt_app = Dispatch("PowerPoint.Application")
    ppt_app.Visible = True  # or False if you prefer it hidden

    # Open the presentation (ReadOnly=True, WithWindow=False)
    pres = ppt_app.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)

    # The SaveAs format 17 corresponds to jpg
    # It will dump slide1.jpg, slide2.jpg, ... into a folder named after the PPT
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    pres.SaveAs(output_dir, 17)

    # Clean up
    pres.Close()
    ppt_app.Quit() 

def create_pptx_with_qr(announcements, output_path, use_summary=False):
    prs = Presentation()

    # Set slide width and height
    prs.slide_width = Inches(13.3333)
    prs.slide_height = Inches(7.5)

    for ann in announcements:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        qr_img = None

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(11.5), Inches(1))
        title_frame = title_box.text_frame

        orig_title = ann['title']
        
        if len(orig_title) > config.TITLE_MAX_CHARS:
            short = summarize_title(orig_title, max_chars=80)
            print(short)
        else:
            short = orig_title

        # # 2) Wrap the (possibly shortened) title into two lines
        # import textwrap
        # wrapped = textwrap.wrap(short, width=config.TITLE_WRAP_WIDTH)
        # if len(wrapped) > 2:
        #     wrapped = wrapped[:2]
        #     wrapped[-1] = wrapped[-1].rstrip(' .,') + "…"

        title_text = short
        title_frame.text = title_text

        title_frame.paragraphs[0].font.size = Pt(42)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*config.BRAND_COLOR_1)  # Dark green
        title_frame.paragraphs[0].font.name = 'Source Sans Pro Black'
        title_frame.paragraphs[0].alignment = 1  # Center alignment
        title_frame.word_wrap = True

        # Body
        if use_summary and 'summary' in ann:
            body_text = ann['summary']

        else:
            body_text = truncate_text(ann['body'], config.MAX_BODY_CHARS)


        font_size = determine_font_size(body_text, config.BASE_FONT_SIZE, config.MIN_FONT_SIZE)
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(2.5))
        body_frame = body_box.text_frame
        body_frame.text = body_text
        body_frame.paragraphs[0].font.size = Pt(font_size)
        body_frame.paragraphs[0].font.name = 'Source Sans Pro'
        body_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black
        body_frame.word_wrap = True

        #Image dynamically scaled
        if ann.get('image_url'):
            img_stream = fetch_image_from_url(ann['image_url'])
            if img_stream:
                # Define the maximum box for the image
                left = Inches(8.53)
                top  = Inches(2.15)
                max_width  = prs.slide_width  - left - Inches(0.5)   # leave 0.5" right margin
                max_height = Inches(5.4)                             # as before

                # Measure the actual image
                img_stream.seek(0)
                with Image.open(img_stream) as img:
                    orig_w, orig_h = img.size
                ratio = orig_w / orig_h

                # Compute scaled dimensions in inches
                if max_width / max_height > ratio:
                    # box is wider than image: scale by height
                    height = max_height
                    width  = height * ratio
                else:
                    # box is taller than image: scale by width
                    width  = max_width
                    height = width / ratio

                # Rewind and insert at computed size
                img_stream.seek(0)
                slide.shapes.add_picture(
                    img_stream,
                    left, top,
                    width=width,
                    height=height
                )

        # QR Code
        if ann.get('link'):
            qr_img = generate_qr_code(ann['link'])
            left = Inches(5.6)
            top = Inches(4.3)
            height = Inches(2.75)
            slide.shapes.add_picture(qr_img, left, top, height=height)

            # Use the actual button text here
            caption = ann.get("button_text", "Scan for more info")
            tb = slide.shapes.add_textbox(Inches(6), Inches(6.95), Inches(2), Inches(0.4))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0,0,0)
            p.alignment = PP_ALIGN.CENTER  # Center alignment
            p.font.italic = True
        else:
            print(f"No link provided for announcement: {ann['title']}")
        
        # Logo
        try:
            slide.shapes.add_picture(config.LOGO_PATH, Inches(0), Inches(6), height=Inches(1.5))
        except FileNotFoundError:
            print(f"Logo file not found at {config.LOGO_PATH}. Skipping logo.")
    prs.save(output_path)
